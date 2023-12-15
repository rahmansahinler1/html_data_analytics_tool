import numpy as np
import pandas as pd
from matplotlib.figure import Figure
from matplotlib.backends.backend_tkagg import (FigureCanvasTkAgg)

import tkinter as tk
from tkinter import *
from tkinter import ttk, filedialog
from tkinter import messagebox, Entry
from tkinter import filedialog

from pptx import Presentation
from pptx.util import Inches

import os
import math
import datetime
from datetime import date, timedelta

class Window(tk.Toplevel):

    def create_plot(self):  
        def open_plot_gui():
            def plot_graph():
                # Figure contains the plot
                figure = Figure(figsize=(16.5, 8.8), dpi=100)
                plot_ = figure.add_subplot(111)

                if not sum([item[1].get() for item in motornummer_wochensoll.values()]):
                    # Raises error if user does not select any motor number and treis to sketch graph
                    messagebox.showerror("Error", f"Please select a motor number!")
                else:
                    # Create canvas to draw
                    canvas = FigureCanvasTkAgg(figure, master=plot_gui)
                    canvas.get_tk_widget().place(x=230, y=6)
                    
                    # Take global cw data
                    global_cw_list = list(self.global_cw_df.values.squeeze(axis=1))
                    initial_year = int(global_cw_list[0].split("-")[0])
                    initial_cw = int(global_cw_list[0].split("-")[1])
                    latest_year = int(global_cw_list[-1].split("-")[0])
                    latest_cw = int(global_cw_list[-1].split("-")[1])
                    
                    # Plot current cw line
                    current_cw = datetime.datetime.now().isocalendar()[1]
                    current_year = datetime.datetime.now().isocalendar()[0]
                    current_cw_index = current_cw - initial_cw + 52 * (current_year - initial_year)
                    plot_.axvline(x=current_cw_index, color='red', linestyle='-')
                    
                    # Comparison between current cw date and latest date
                    if current_year > latest_year:
                        latest_year = current_year
                        latest_cw = current_cw
                    elif current_year == latest_year:
                        latest_cw = max(latest_cw, current_cw)
                   
                    # Compare the milestone dates with the global_cw_list, hold max and min y value for later scatter plot
                    milestone_dates = {}
                    milestone_indexes = []
                    for index, (key, value) in enumerate(dict_milestone.items()):
                        milestone_name = key
                        milestone_date = value
                        milestone_year = int(milestone_date.split("/")[2])
                        milestone_cw = datetime.datetime.strptime(milestone_date, '%m/%d/%Y').isocalendar()[1]
                        milestone_index = milestone_cw - initial_cw
                        if milestone_year > initial_year: milestone_index += (milestone_year - initial_year) * 52
                        # Do we sketch corresponding milestone?
                        if milestone_selection[index].get(): plot_.axvline(x=milestone_index, color='black', linestyle='-')
                        milestone_indexes.append(milestone_index)   
                        # Comparison between plot data and milestone to find latest date"
                        if milestone_year > latest_year:
                            latest_year = milestone_year
                            latest_cw = milestone_cw
                        elif milestone_year == latest_year:
                            latest_cw = max(latest_cw, milestone_cw)                 
                        # Append milestone dates for later comparison
                        milestone_dates[milestone_name] = [milestone_year, milestone_cw]                
                    
                    x_length = 0        # Variable to hold length of x-axis
                    color_index = 0     # Variable to hold and reset color selection
                    max_y_value = 0     # Variable to hold maximum y axis value for scatter plot
                    for motornummer, value in motornummer_wochensoll.items():
                        if sum([item[1].get() for item in motornummer_wochensoll.values()]) > 10:
                            messagebox.showerror("Error", f"Up to 10 variants can be selected at the same time!")
                            break
                        elif value[1].get():
                            try:
                                step_difference = float(value[0].get())     # km / week value in the entry box
                            except ValueError:
                                messagebox.showerror("Error", f"Please enter only numerical values for Wochensoll!")
                                break     
                            # Filter corresponding motor numbers
                            group = self.db_df[self.db_df['Motornummer'] == motornummer]
                            
                            # PST and Versuchsbeschreibung value, 'first value for corresponding motor number'
                            pst = group['PST'].iloc[0]
                            versuchsbeschreibung = group['Versuchsbeschreibung'].iloc[0]
                            
                            # Legend: {Motornumber}\n{PST}/{Versuchsbeschreibung}
                            legend_label_str = f"{motornummer}\n{pst} / {versuchsbeschreibung}"                                                                                                           
                            
                            # Index calculation: Zeitstempel to Index
                            group_indexes = []      # [year, cw]
                            for _, row in group.iterrows():
                                row_year = int(row['Zeitstempel'].split('-')[0])
                                row_cw = datetime.datetime.strptime(row['Zeitstempel'], '%Y-%m-%d').isocalendar()[1]
                                row_index = row_cw - initial_cw
                                if row_year > initial_year: row_index += (row_year - initial_year) * 52
                                if row_index < 0: row_index = 0
                                group_indexes.append(row_index)
                            
                            # If reached aimed value, do not increase the size else increase it with latest value to the current cw
                            actual_laufstrecke = group['Laufstrecke'].values[~np.isnan(group['Laufstrecke'].values)]
                            actual_sollaufstrecke = group['Solllaufstrecke'].values[~np.isnan(group['Solllaufstrecke'].values)]
                            first_laustrecke, last_laufstrecke, last_sollaufstrecke = actual_laufstrecke[0], actual_laufstrecke[-1], actual_sollaufstrecke[-1]
                            if last_laufstrecke >= last_sollaufstrecke:
                                # Laufstrecke
                                plot_.plot(group_indexes, group['Laufstrecke'].values, '-o', label=legend_label_str,
                                       color=self.colordb[color_index])
                                
                                # Solllaufstrecke
                                soll_line_step = math.ceil((last_sollaufstrecke - first_laustrecke) / step_difference)
                                soll_line_values = np.linspace(group['Laufstrecke'].values[0], group['Solllaufstrecke'].values[-1], soll_line_step + 1)
                                soll_line_index = pd.RangeIndex(group_indexes[0], group_indexes[0] + len(soll_line_values))
                                plot_.plot(soll_line_index, soll_line_values, '-.',
                                       color=self.colordb[color_index])
                            else:
                                # Extend graphic line with the planned way
                                ext_line_step = math.ceil((last_sollaufstrecke - last_laufstrecke) / step_difference)
                                ext_line_values = np.linspace(group['Laufstrecke'].values[-1], group['Solllaufstrecke'].values[-1], ext_line_step + 1)
                                ext_line_index = pd.RangeIndex(current_cw_index, current_cw_index + len(ext_line_values))
                                x_length = max(x_length, ext_line_index[-1])
                                plot_.plot(ext_line_index, ext_line_values, '-o', lw=2.0,
                                    color=self.colordb_faded[color_index])
                                
                                # Laufstrecke line
                                index = pd.RangeIndex(group_indexes[0], current_cw_index + 1)
                                diff = len(index) - len(group['Laufstrecke'].values)                                        
                                padded_laufstrecke = np.pad(group['Laufstrecke'].values, (0, diff), 'constant', constant_values=last_laufstrecke)
                                plot_.plot(index, padded_laufstrecke, '-o', label=legend_label_str,
                                       color=self.colordb[color_index])
                                
                                # Solllaufstrecke
                                soll_line_step = math.ceil((last_sollaufstrecke - first_laustrecke) / step_difference)
                                soll_line_values = np.linspace(group['Laufstrecke'].values[0], group['Solllaufstrecke'].values[-1], soll_line_step + 1)
                                soll_line_index = pd.RangeIndex(group_indexes[0], group_indexes[0] + len(soll_line_values))
                                plot_.plot(soll_line_index, soll_line_values, '-.',
                                       color=self.colordb[color_index])
                                
                            # Increase color index for later selections
                            color_index += 1  
                            # Take current maximum value for y axis
                            max_y_value = max(max_y_value, np.nanmax(group['Solllaufstrecke'].values))                                       
                            
                    # Scatter plot for markers of milestones
                    markers = ['o','v','s','P','*','X','D','h','+','x']
                    y_values = np.linspace(0, max_y_value, 10)
                    for j, milestone_index in enumerate(milestone_indexes):
                        if milestone_selection[j].get():
                            # Do we sketch corresponding milestone?
                            plot_.scatter([milestone_index] * len(y_values), y_values, marker=markers[j], color='black')         
                    # Create a date list start from inital_year, initial_week to latest_year, latest_week. 
                    # Format: year- .week. If Milestone: milestone_name year- .week
                    x_formatted = []
                    milestone_x_length = latest_cw - initial_cw
                    if latest_year > initial_year:
                        milestone_x_length += (latest_year - initial_year) * 52
                    x_length = max(x_length, milestone_x_length)                        
                    
                    year = initial_year
                    week = initial_cw
                    x_length += 3
                    for _ in range(x_length + 1):
                        year_change_flag = False
                        milestone_change_flag = False     
                        # Year change detection and string value creation
                        if week > 52:
                            week = 1
                            year += 1
                            year_change_flag = True 
                        if week < 10:
                            week_str = '0' + str(week)
                        else:
                            week_str = str(week)
                        year_str = str(year)
                        
                        # Compare milestone date and index date to change x-axis text
                        milestone_name = ''
                        for index, (key, value) in enumerate(milestone_dates.items()):
                            if ([year, week] == value) and (milestone_selection[index].get()):
                                # If milestone is selected and indexes [year, week] matches
                                milestone_name = key
                                milestone_change_flag = True              
                        
                        # If first and every year change, enter full year
                        formatted_date = ''
                        if week == initial_cw and year == initial_year:
                            formatted_date = (f'{year_str}-.{week_str}')
                        elif year_change_flag:
                            formatted_date = (f'{year_str}-.{week_str}')
                        else:
                            formatted_date = (f'-.{week_str}')
                        
                        if milestone_change_flag:
                            formatted_date = f'{milestone_name} {formatted_date}'
                        
                        x_formatted.append(formatted_date)
                        week += 1
                          
                    plot_.set_autoscalex_on(False)
                    plot_.set_xticks(np.arange(len(x_formatted)))
                    plot_.set_xticklabels(x_formatted, rotation=90, fontsize=9)          
                    plot_.set_title("Zyklus")
                    plot_.set_ylabel('Laufstrecke und Solllaufstrecke [Km]')
                    plot_.set_xlabel('Kalenderwoche', fontsize=10)
                    xmin, xmax = plot_.get_xlim()
                    plot_.set_xlim(xmin - 0.5, xmax)
                    plot_.grid(True, axis='y')
                    legend_ = plot_.legend(loc='upper left', framealpha=0.5)
                    legend_.set_draggable(state=True)
                    self.plot_gui_custom_figure = figure
                    figure.tight_layout()
                    canvas.draw()

            def save_graph():
                # Ask the user for the save file
                plot_path = filedialog.asksaveasfilename(
                    filetypes=[("png file", ".png")],
                    defaultextension=".png",
                    initialdir=os.getcwd())

                if len(plot_path) == 0:
                    return

                self.plot_gui_custom_figure.savefig(plot_path, dpi=100)

                # Create message box
                messagebox.showinfo("Info", f"Plot file is saved!\n\nFile Path: {plot_path}")

            # Toplevel object which will be treated as a new window
            plot_gui = Toplevel()
            plot_gui.title("Sketch Window")
            plot_gui.resizable(False, False)

            # GUI Geometry
            plot_gui.geometry('1900x900+5+50')  # adjusts geometry geometry(widthxheight+x+y) x,y: offset in pixels
            proj_dir = os.getcwd()
            plot_gui.iconbitmap(f'{proj_dir}\\utils\\avl_logo.ico')  # changes icon of the windows
            
            # Motor Number selection and Wochensoll values
            label = Label(plot_gui, text='Motornummer - \u0394Km / Woche')
            label.config(font=("Arial", 9), fg="white", bg="#7F7F7F")
            label.place(x=5, y=5, width=200, height=15)
            
            motornummer = sorted(self.db_df_uniques_keep_first["Motornummer"].to_list())
            motornummer_wochensoll = {key:['', 0] for key in motornummer}
            motornummer_selection = [tk.IntVar() for _ in range(len(motornummer))]
            latest_widget_place = 0
            
            # Motor number labels, checkboxes and km/week value entrybox
            for index, item in enumerate(motornummer):
                motornummer_frame = Frame(plot_gui)
                motornummer_frame_index = 25 + (index * 20)
                motornummer_frame.place(x=5, y=motornummer_frame_index, width=200, height=20)
                
                # Checkbox for motornumber selection
                checkbox = Checkbutton(motornummer_frame, text='', variable=motornummer_selection[index], anchor="w")
                checkbox.pack(side=LEFT)
                
                # Label for the motor number
                label = Label(motornummer_frame, text=item)
                label.pack(side=LEFT)
                
                # Entry for user input
                wohcensoll_var = tk.StringVar(value='7500')
                entry = Entry(motornummer_frame, width=10, textvariable=wohcensoll_var)
                entry.pack(side=RIGHT) 
                
                # Store the entry in the dictionary
                motornummer_wochensoll[item] = [entry, motornummer_selection[index]]
                
                # Update latest motornummer place
                latest_widget_place = motornummer_frame_index
 
            # Milestone Section
            dict_milestone = self.get_milestone()
            
            # Milestone selection frames
            if dict_milestone:
                # Milestone label
                label2 = Label(plot_gui, text='Milestone Names & Selection')
                label2.config(font=("Arial", 9), fg="white", bg="#7F7F7F")
                label2.place(x=5, y=latest_widget_place + 25, width=200, height=15)

                # Selection holder value
                milestone_selection = [tk.IntVar() for _ in range(len(dict_milestone))]     
                
                # Milestone selection frames
                for index, (key, value) in enumerate(dict_milestone.items()):
                    milestone_frame = Frame(plot_gui, relief=FLAT)
                    milestone_frame_index = latest_widget_place + 45 + (index * 20)
                    milestone_frame.place(x=5, y=milestone_frame_index, width=200, height=20)

                    # Checkbox for motornumber selection
                    checkbox = Checkbutton(milestone_frame, text='', variable=milestone_selection[index], anchor="w")
                    checkbox.pack(side=LEFT)
                    
                    # Label for the milestone name
                    label = Label(milestone_frame, text=key)
                    label.pack(side=LEFT)
                    
                    # Label for the milestone date
                    label1 = Label(milestone_frame, text=value)
                    label1.pack(side=RIGHT)
                    
                # Save latest milestone place
                latest_widget_place = milestone_frame_index
            
            # Plot button
            plot_button = ttk.Button(plot_gui, text="Plot", command=plot_graph)
            plot_button.place(x=5, y=latest_widget_place + 30, width=90, height=25)

            # Save figure button
            save_button = ttk.Button(plot_gui, text="Save", command=save_graph)
            save_button.place(x=115, y=latest_widget_place + 30, width=90, height=25)
                    
            # Information for wochensoll value enterance
            label1 = Label(plot_gui, text='!Only numerical values for \u0394Km / Woche', anchor='w')
            label1.config(font=("Arial", 7, 'italic', 'bold'), fg="black")
            label1.place(x=5, y=latest_widget_place + 60, width=200, height=15)
            
            # Information for milestone selection value enterance
            if dict_milestone:
                label2 = Label(plot_gui, text='Select checkbox to sketch milestone line', anchor='w')
                label2.config(font=("Arial", 7, 'italic', 'bold'), fg="black")
                label2.place(x=5, y=latest_widget_place + 77, width=200, height=15)
                
            # Current CW Label
            label3 = Label(plot_gui, text=f'Current CW: {datetime.datetime.now().isocalendar()[1]}')
            label3.config(font=("Arial", 10), fg="white", bg="#7F7F7F")
            label3.place(x=5, y=880, width=200, height=15)

        # Color Map
        self.colordb = ['#0072BD', '#D95319', '#EDB120', '#7E2F8E', '#77AC30', '#4DBEEE', '#A2142F', '#DA0017', '#606060', '#F267B2']
        self.colordb_faded = ['#80B9DE', '#ECA69D', '#F6D18C', '#BF90C7', '#BBD966', '#A6DFEF', '#D15777', '#ED003B', '#A0A0A0', '#F984C6']
        self.colordb_number = len(self.colordb)

        # Update user
        messagebox.showinfo("Info", f"Additional GUI oppenning for sketching plot")
        open_plot_gui()

    def export_xlsx(self):
        xlsx_path = filedialog.asksaveasfilename(
            filetypes=[("Excel file", ".xlsx")],
            defaultextension=".xlsx",
            initialdir=os.getcwd())
        if len(xlsx_path) == 0:
            return
        else:
            try:
                # Create a Pandas Excel writer using XlsxWriter as the engine.
                writer = pd.ExcelWriter(xlsx_path, engine='xlsxwriter')
            except:
                messagebox.showerror("Error", f"Please ensure that the file is not open by an external application!")
                return

        # export df to excel
        table_df = self.df_table
        table_df.to_excel(writer, sheet_name='Table', startrow=1, index=False, header=False)

        # Get the xlsxwriter workbook and worksheet objects.
        workbook = writer.book
        worksheet = writer.sheets['Table']

        # Add a header format.
        header_format_grey = workbook.add_format({
            'bold': True,
            'text_wrap': True,
            'valign': 'top',
            'fg_color': '#c8cacf',      # grey
            'border': 1})

        header_format_yellow = workbook.add_format({
            'bold': True,
            'text_wrap': True,
            'valign': 'top',
            'fg_color': '#f1f505',      # yellow
            'border': 1})

        header_format_orange = workbook.add_format({
            'bold': True,
            'text_wrap': True,
            'valign': 'top',
            'fg_color': '#f58402',      # orange
            'border': 1})

        # Write the column headers with the defined format.
        for col_num, value in enumerate(table_df.columns.values):
            if col_num == 0:
                worksheet.write(0, col_num, value, header_format_grey)
            elif col_num % 8 >= 1 and col_num % 8 <= 4:
                worksheet.write(0, col_num, value, header_format_yellow)
            else:
                worksheet.write(0, col_num, value, header_format_orange)

        writer.close()
        messagebox.showinfo("Info", f"Excel file is created!\n\nFile Path: {xlsx_path}")

    def export_pptx(self):
        def open_select_files_gui():
            def raise_above_all(window):
                # bring front after adding, deleting and saving file process
                window.attributes('-topmost', 1)
                window.attributes('-topmost', 0)

            def add_files():
                # Ask the user for the save file
                file_type = (('PNG files', '*.png'), ('JPG files', '*.jpg'))
                image_files = filedialog.askopenfilenames(title='Select files', initialdir=os.getcwd(), filetypes=file_type)

                if len(image_files) == 0:
                    return
                else:
                    for each_image_path in image_files:
                        each_file_name = each_image_path.split("/")[-1]
                        self.export_image_file_list.append((each_image_path, each_file_name))

                    list_box.delete(0, END)
                    # Show html file names on the tool
                    for ctr_file, (_, name) in enumerate(self.export_image_file_list):
                        list_box.insert(ctr_file, name)
                        
                    raise_above_all(plot_gui)

            def delete_selected_files():
                for ctr_file, (_, _) in enumerate(self.export_image_file_list):
                    list_box.delete(ctr_file)
                    self.export_image_file_list.pop(ctr_file)

                raise_above_all(plot_gui)

            def move_slide(xml_slides, old_index, new_index):
                slides = list(xml_slides)
                xml_slides.remove(slides[old_index])
                xml_slides.insert(new_index, slides[old_index])

            def save_image_files():
                proj_dir = os.getcwd()
                prs = Presentation(f"{proj_dir}\\utils\\Template.pptx")
                layout = prs.slide_layouts[5]

                for ctr_slide, (img_path, f_name) in enumerate(self.export_image_file_list):
                    slide = prs.slides.add_slide(layout)
                    xml_slides = prs.slides._sldIdLst
                    move_slide(xml_slides, -1, ctr_slide+1)  # -1 for last slide position after addition. 1: means the second slide position in the presentation
                    title = slide.shapes.title
                    title.text = f_name.rsplit('.', 1)[0]
                    slide.shapes.add_picture(img_path, Inches(0.6), Inches(1.4), height=Inches(5.6))

                # Save presentation
                presentation_path = filedialog.asksaveasfilename(
                    filetypes=[("Powerpoint File", ".pptx")],
                    defaultextension=".pptx",
                    initialdir=os.getcwd())

                if len(presentation_path) == 0:  # cancel button is pressed
                    raise_above_all(plot_gui)
                    return
                else:
                    try:
                        prs.save(presentation_path)
                        messagebox.showinfo('Info', f'Powerpoint file created!\n\n {presentation_path}')
                        raise_above_all(plot_gui)
                    except:
                        messagebox.showerror("Error", f"Please ensure that the file is not open by an external application!")
                        raise_above_all(plot_gui)
                        return

            # Toplevel object which will be treated as a new window
            plot_gui = Toplevel()
            plot_gui.title("Export PPT GUI")
            plot_gui.resizable(False, False)

            # GUI Geometry
            plot_gui.geometry('255x220+830+250')  # adjusts geometry geometry(widthxheight+x+y) x,y: offset in pixels
            proj_dir = os.getcwd()

            # Select file button
            select_button = ttk.Button(plot_gui, text="Add File", command=add_files)
            select_button.place(x=5, y=180)

            # Delete file button
            delete_button = ttk.Button(plot_gui, text="Delete File", command=delete_selected_files)
            delete_button.place(x=90, y=180)

            # Save files button
            save_button = ttk.Button(plot_gui, text="Save", command=save_image_files)
            save_button.place(x=175, y=180)

            # List Box
            list_box = Listbox(plot_gui, selectmode=SINGLE, width=40, height=10)
            list_box.place(x=5, y=5)

        self.export_image_file_list = []
        open_select_files_gui()

    def create_plan(self):

        def get_end_date(group, global_start_date):

            distance_per_week = 7500.0  # must be integer for range usage

            last_actual_distance = group.loc[len(group)-1, "Laufstrecke"]
            last_target_distance = group.loc[len(group)-1, "Solllaufstrecke"]
            distance_difference = last_target_distance - last_actual_distance

            num_travel_week = distance_difference / distance_per_week
            num_travel_day = num_travel_week * 7.0

            if last_actual_distance >= (last_target_distance * 0.99): # the desired distance
                plan_end_datetime = pd.to_datetime(group.loc[len(group)-1, "Zeitstempel"])
            else:
                # Find next monday to start adding extra travel distance duration on that date
                today = date.today()
                start_of_week = today - timedelta(days=today.weekday())
                day_diff = datetime.datetime.today().weekday() + pd.to_datetime(global_start_date).weekday() - datetime.datetime.today().weekday()
                plan_end_datetime = start_of_week + timedelta(days= (num_travel_day + day_diff))

            if plan_end_datetime.day < 10:
                plan_end_date_day = "0" + str(plan_end_datetime.day)
            else:
                plan_end_date_day = str(plan_end_datetime.day)

            if plan_end_datetime.month < 10:
                plan_end_date_month = "0" + str(plan_end_datetime.month)
            else:
                plan_end_date_month = str(plan_end_datetime.month)

            plan_end_date = str(plan_end_datetime.year) + "-" + plan_end_date_month + "-" + plan_end_date_day
            group.loc[:, "Plan_end_date"] = plan_end_date

            return group

        def detect_stop_duration():
            df_create_plan_stop_list = []
            group_list_last = self.df_plan_list_last_values.copy()
            # Loop searches through individual motornumber data and creates a list for them to store if there is input for that corresponding cw
            for group in group_list_last:
                group_stop_list = [0] * len(group)
                for row_index, row in group.iterrows():
                    try:
                        int(row['Laufstrecke'])
                        group_stop_list[row_index] = 1
                    except ValueError:
                        group_stop_list[row_index] = 0     
                df_create_plan_stop_list.append(group_stop_list)
            return df_create_plan_stop_list

        def open_plan_gui():
            def plot_plan():
                if sum([item.get() for item in motornummer_selection.values()]) < 2:
                    # Raises error if user does not select any motor number and treis to sketch graph
                    messagebox.showerror("Error", f"Please select at least 2 motor numbers to sketch plan")
                else:
                    # Get the selected motor numbers
                    motor_numbers = [] 
                    for motor_number, selection in motornummer_selection.items():
                        if selection.get(): motor_numbers.append(motor_number)
                    
                    txt_ch = (self.db_df[self.db_df['Motornummer'].isin(motor_numbers)]
                            .sort_values(by="Start_date", ascending=True)
                            .reset_index(drop=True))
                    
                    # Figure contains the plot
                    figure = Figure(figsize=(16.5, 8.8), dpi=100)
                    axs = figure.subplots(len(txt_ch['Motornummer'].unique()), sharex=True)
                    figure.subplots_adjust(hspace=0)
                    
                    # Create canvas to draw
                    canvas = FigureCanvasTkAgg(figure, master=plan_gui)
                    canvas.get_tk_widget().place(x=230, y=6)
                    
                    # Time and Dates   
                    current_time = datetime.datetime.now().timestamp()
                    current_date = datetime.date.fromtimestamp(current_time).isocalendar()
                    current_year = current_date[0]
                    current_cw = current_date[1]
                    hold_max_milestone_year_week = [0, 0]
                    initial_cw = int(txt_ch.reset_index(drop=True).loc[0, "KW"].split("-")[1])  # Database initial week
                    initial_year = int(txt_ch.reset_index(drop=True).loc[0, "KW"].split("-")[0])  # Database initial year
                    
                    # Markers for milestone lines
                    markers = ['o','v','s','P','*','X','D','h','+','x']
                    
                    # Select the place for current CW line, cover the case that year passed
                    relative_current_cw = current_cw + 52 * (current_year - initial_year) - initial_cw
                    
                    # Create plot for each motor number
                    for i in range(len(txt_ch['Motornummer'].unique())):
                        # Create CW plot layout
                        axs[i].set(yticklabels=[])
                        axs[i].set_ylabel(sorted(txt_ch['Motornummer'].unique())[i], rotation="horizontal")
                        axs[i].yaxis.set_label_coords(-0.05, 0.5)
                        axs[i].tick_params(left=False)
                        axs[i].axvline(relative_current_cw, label="aktuelle KW", color="r")
                    
                        # Add milestone lines
                        for j, child in enumerate(self.milestone_tree.get_children()):
                            if milestone_selection[j].get():
                                milestone_name = self.milestone_tree.item(child)['values'][0]
                                milestone_date = self.milestone_tree.item(child)['values'][1]
                                splitted_milestone_date = milestone_date.split("/")
                                milestone_week = datetime.date(int(splitted_milestone_date[2]), int(splitted_milestone_date[0]), int(splitted_milestone_date[1])).isocalendar()[1]
                                milestone_year = int(splitted_milestone_date[2])
                                
                                # Holding maximum milestone week for following comparsion
                                if milestone_year > hold_max_milestone_year_week[0]:
                                        hold_max_milestone_year_week = [milestone_year, milestone_week]
                                elif milestone_year == hold_max_milestone_year_week[0] and \
                                    milestone_week > hold_max_milestone_year_week[1]:
                                        hold_max_milestone_year_week = [milestone_year, milestone_week]              
                                
                                # Selecting place to where to add indicator of CW to the plan, cover the case that year passed for first database entry
                                relative_milestone_week = milestone_week + 52 * (milestone_year - initial_year) - initial_cw
                                axs[i].axvline(relative_milestone_week, label="# " + milestone_name, linestyle='-', marker=markers[j], color="k")    
                        axs[i].grid(color='lightgrey', linestyle='--', linewidth=0.5, axis='x')
                    
                    handles, labels = axs[0].get_legend_handles_labels()
                    
                    create_plan_stop_list = detect_stop_duration()  
                    global_start_time = pd.to_datetime(txt_ch['Start_date']).iloc[0].timestamp() 
                    global_start_date = datetime.datetime.fromtimestamp(global_start_time).isocalendar()
                    start_first_global = global_start_time / 3600.0 / 24.0 / 7.0
                    
                    hold_max_end_plan_year_week = [0, 0]
                    
                    arrow_colors = {0: ['#E8E8E8', 'Standzeit'], 
                            1: ['#00437A', 'Ist-Laufzeit'],  
                            2: ['#00437A80', 'Soll-Laufzeit']}  # 0: light-grey, 1: blue, 2: light-blue
                    arrow_counter = list(arrow_colors.keys())
                    
                    for i, group in enumerate(txt_ch.groupby('Motornummer')):

                        group = group[1].sort_values(by='Zeitstempel',ascending=True)
                        group = group.reset_index(drop=True)
                        group = get_end_date(group, txt_ch['Start_date'].iloc[0])

                        start_day_diff = pd.to_datetime(group['Start_date']).iloc[0].weekday() - pd.to_datetime(txt_ch['Start_date']).iloc[0].weekday()
                        start_first_group = (pd.to_datetime(group['Start_date']).iloc[0].timestamp() / 3600.0 / 24.0 - start_day_diff) / 7.0
                        end_plan_time = pd.to_datetime(group['Plan_end_date']).iloc[-1].timestamp()
                        end_plan_date = datetime.datetime.fromtimestamp(end_plan_time).isocalendar()
                        end_plan_year_week = [end_plan_date[0], end_plan_date[1]]
                        total_distance = group['Solllaufstrecke'].iloc[-1]
                        completed_distance = group['Laufstrecke'].iloc[-1]
                                
                        # Holding maximum end plan week for following comparsion
                        if end_plan_year_week[0] > hold_max_end_plan_year_week[0]:
                                hold_max_end_plan_year_week = [end_plan_year_week[0], end_plan_year_week[1]]
                        elif end_plan_year_week[0] == hold_max_end_plan_year_week[0] and \
                            end_plan_year_week[1] > hold_max_end_plan_year_week[1]:
                                    hold_max_end_plan_year_week = [end_plan_year_week[0], end_plan_year_week[1]]
                                                        
                        # Plotting the plan
                        end_plan_group = end_plan_time / 3600.0 / 24.0 / 7.0
                        diff = end_plan_group - start_first_group
                        stop_duration = create_plan_stop_list[i]
                
                        if completed_distance >= total_distance * 0.99:  # 99% of the desired distance
                            arrow_stop = axs[i].arrow((start_first_group-start_first_global), 0, diff, 0, width=0.02,
                                            length_includes_head=True, head_width=0.05275,
                                            head_length=0.2, color=arrow_colors[1][0])
                            
                        else:
                            skip_at_start_counter = next((i for i, x in enumerate(stop_duration) if x), None) # x!= 0 for strict match
                            if skip_at_start_counter:
                                diff += skip_at_start_counter  # to make the length of the list equal to the diff
                            
                            # Fill stop druation with 0s equal to length of the diff
                            stop_duration = stop_duration + [0] * (math.ceil(diff) - len(stop_duration))  # to make the length of the list equal to the diff
                            arrows_after_ccw = math.ceil(diff - (current_cw - initial_cw + 52 * (current_year - initial_year)))  # additional arrows to be added to the plot after current cw
                            if arrows_after_ccw:
                                stop_duration[-arrows_after_ccw:] = [2] * (arrows_after_ccw)
                            
                            for index, arrow_type in enumerate(stop_duration):
                                if skip_at_start_counter:
                                    skip_at_start_counter -= 1
                                    continue
                                
                                if index == len(stop_duration) - 1:
                                    head_length = 0.2650
                                    head_width = 0.05275
                                else:
                                    head_length = 0
                                    head_width = 0
                                    
                                arrow_travel = axs[i].arrow(index, 0,
                                            1, 0, width=0.02, length_includes_head=True, head_width=head_width,
                                            head_length=head_length, color=arrow_colors[arrow_type][0])
                                    
                                if arrow_type in arrow_counter:
                                    arrow_counter.remove(arrow_type)
                                    handles.append(arrow_travel)
                                    labels.append(arrow_colors[arrow_type][1])
                    
                            # Add left distance text into the plot
                            axs[i].text(len(stop_duration) + 0.5, 0,
                                        f'{int(completed_distance)} Km ({int(total_distance)} tKm)',
                                        fontsize=10)
                            
                    # Comparing latest milestone date and latest shot date to select
                    if hold_max_milestone_year_week[0] > hold_max_end_plan_year_week[0]:
                        if hold_max_milestone_year_week[0] > initial_year:
                            max_week = hold_max_milestone_year_week[1] + 52 *(hold_max_milestone_year_week[0] - initial_year)  - global_start_date[1]
                        else:
                            max_week = hold_max_milestone_year_week[1] - global_start_date[1]
                    elif hold_max_end_plan_year_week[0] > hold_max_milestone_year_week[0]:
                        if hold_max_end_plan_year_week[0] > initial_year:
                            max_week = hold_max_end_plan_year_week[1] + 52 *(hold_max_end_plan_year_week[0] - initial_year)  - global_start_date[1]
                        else:
                            max_week = hold_max_end_plan_year_week[1] - global_start_date[1]
                    else:
                        if hold_max_milestone_year_week[1] > hold_max_end_plan_year_week[1]:
                            if hold_max_milestone_year_week[0] > initial_year:
                                max_week = hold_max_milestone_year_week[1] + 52 *(hold_max_milestone_year_week[0] - initial_year)  - global_start_date[1]
                            else:
                                max_week = hold_max_milestone_year_week[1] - global_start_date[1]
                        else:
                            if hold_max_end_plan_year_week[0] > initial_year:
                                max_week = hold_max_end_plan_year_week[1] + 52 *(hold_max_end_plan_year_week[0] - initial_year)  - global_start_date[1]
                            else:
                                max_week = hold_max_end_plan_year_week[1] - global_start_date[1]

                    # Adjust x-axis of the CW plot
                    calendar_weeks = []
                    hold_year = initial_year
                    hold_week = initial_cw
                    # Create a list of KWs starting from global initial to global end
                    for week_counter in range(hold_week, hold_week + max_week + 1):
                        if week_counter == hold_week:
                            calendar_weeks.append(f'{hold_year}- .{hold_week}')
                        elif week_counter % 52 == 0:
                            hold_year += 1
                            hold_week = 1
                            calendar_weeks.append(f'{hold_year}- .{hold_week}')
                        else:
                            hold_week += 1
                            calendar_weeks.append(f'.{hold_week}')
                            
                    figure.gca().set_autoscalex_on(False)
                    figure.suptitle('Zeitplan', fontsize=14)
                    figure.text(0.04, 0.5, 'Motornummer', va='center', rotation='vertical', fontsize=12)
                    figure.text(0.5, 0.04, 'Kalenderwoche', ha='center', fontsize=12)
                    figure.gca().set_xticks(np.arange(max_week+1))
                    figure.gca().set_xticklabels(calendar_weeks, rotation=90, fontsize=9)
                    legend_ = figure.gca().legend(loc='lower left', framealpha=0.5, 
                                        handles=handles, labels=labels,
                                        bbox_to_anchor=(0.96, 0))
                    legend_.set_draggable(state=True)
                    xmin, xmax = figure.gca().get_xlim()
                    figure.gca().set_xlim(xmin - 0.5, xmax)
                    self.plan_gui_custom_figure = figure
                    canvas.draw()
                                       
            def save_plan():
                # Ask the user for the save file
                plan_path = filedialog.asksaveasfilename(
                    filetypes=[("png file", ".png")],
                    defaultextension=".png",
                    initialdir=os.getcwd())
                if not plan_path: return
                self.plan_gui_custom_figure.savefig(plan_path, dpi=100)
                # Create message box
                messagebox.showinfo("Info", f"Plan Graphic is saved!\n\nFile Path: {plan_path}")

            def select_all_motornumbers():
                for item in motornummer_selection_values: item.set(1)
            
            def deselect_all_motornumbers():
                for item in motornummer_selection_values: item.set(0)
                
            # Toplevel object which will be treated as a new window
            plan_gui = Toplevel()
            plan_gui.title("Sketch Window")
            plan_gui.resizable(False, False)
            
            # GUI Geometry
            plan_gui.geometry('1900x900+5+50')  # adjusts geometry geometry(widthxheight+x+y) x,y: offset in pixels
            proj_dir = os.getcwd()
            
            # Motor Number selection and Wochensoll values
            label = Label(plan_gui, text='Motornumber Selection')
            label.config(font=("Arial", 9), fg="white", bg="#7F7F7F")
            label.place(x=5, y=5, width=200, height=15)
            
            motornummer = sorted(self.db_df_uniques_keep_first["Motornummer"].to_list())
            motornummer_selection = {key:1 for key in motornummer}
            motornummer_selection_values = [tk.IntVar(value=1) for _ in range(len(motornummer))]
            latest_widget_place = 0
            
            # Motor number labels, checkboxes and km/week value entrybox
            for index, item in enumerate(motornummer):
                motornummer_frame = Frame(plan_gui)
                motornummer_frame_index = 25 + (index * 20)
                motornummer_frame.place(x=5, y=motornummer_frame_index, width=200, height=20)
                
                # Checkbox for motornumber selection
                checkbox = Checkbutton(motornummer_frame, text='', variable=motornummer_selection_values[index], anchor="w")
                checkbox.pack(side=LEFT)
                
                # Label for the motor number
                label = Label(motornummer_frame, text=item, )
                label.pack(expand=True)
                
                # Store the entry in the dictionary
                motornummer_selection[item] = motornummer_selection_values[index]
                
                # Update latest motornummer place
                latest_widget_place = motornummer_frame_index
            
            # Select all / Deselect all buttons
            select_all_button = ttk.Button(plan_gui, text="Select All", command=select_all_motornumbers)
            select_all_button.place(x=5, y=latest_widget_place + 25, width=90, height=25)
            
            # Select all / Deselect all buttons
            deselect_all_button = ttk.Button(plan_gui, text="Deselect All", command=deselect_all_motornumbers)
            deselect_all_button.place(x=115, y=latest_widget_place + 25, width=90, height=25)
            
            # Milestone Section
            dict_milestone = self.get_milestone()
            
            # Milestone selection frames
            if dict_milestone:
                # Milestone label
                label2 = Label(plan_gui, text='Milestone Names & Selection')
                label2.config(font=("Arial", 9), fg="white", bg="#7F7F7F")
                label2.place(x=5, y=latest_widget_place + 60, width=200, height=15)

                # Selection holder value
                milestone_selection = [tk.IntVar() for _ in range(len(dict_milestone))]     
                
                # Milestone selection frames
                for index, (key, value) in enumerate(dict_milestone.items()):
                    milestone_frame = Frame(plan_gui, relief=FLAT)
                    milestone_frame_index = latest_widget_place + 80 + (index * 20)
                    milestone_frame.place(x=5, y=milestone_frame_index, width=200, height=20)

                    # Checkbox for motornumber selection
                    checkbox = Checkbutton(milestone_frame, text='', variable=milestone_selection[index], anchor="w")
                    checkbox.pack(side=LEFT)
                    
                    # Label for the milestone name
                    label = Label(milestone_frame, text=key)
                    label.pack(side=LEFT)
                    
                    # Label for the milestone date
                    label1 = Label(milestone_frame, text=value)
                    label1.pack(side=RIGHT)
                    
                # Save latest milestone place
                latest_widget_place = milestone_frame_index
            
            # Plot button
            plot_button = ttk.Button(plan_gui, text="Plot", command=plot_plan)
            plot_button.place(x=5, y=latest_widget_place + 30, width=90, height=25)
            
            # Save figure button
            save_button = ttk.Button(plan_gui, text="Save", command=save_plan)
            save_button.place(x=115, y=latest_widget_place + 30, width=90, height=25)
            
            # Information for milestone selection value enterance
            if dict_milestone:
                label2 = Label(plan_gui, text='Select checkbox to sketch', anchor='w')
                label2.config(font=("Arial", 7, 'italic', 'bold'), fg="black")
                label2.place(x=5, y=latest_widget_place + 58, width=200, height=15)
            
            # Current CW Label
            label3 = Label(plan_gui, text=f'Current CW: {datetime.datetime.now().isocalendar()[1]}')
            label3.config(font=("Arial", 10), fg="white", bg="#7F7F7F")
            label3.place(x=5, y=880, width=200, height=15)
        
        # Update user
        messagebox.showinfo("Info", f"Additional GUI oppenning for sketching plot")
        open_plan_gui()
    
    def select_date(self):
        selected_date = pd.to_datetime(self.calendar.get_date()).strftime('%m/%d/%Y')
        selected_milestone_name = self.get_milestone_name()
        self.milestone_tree.insert('', tk.END, values=[selected_milestone_name, selected_date])
    
    def update_memory(self):
        dict_milestone = {'Name':[], 'Date':[]}
        for child in self.milestone_tree.get_children():
            milestone_name = self.milestone_tree.item(child)['values'][0]
            milestone_date = self.milestone_tree.item(child)['values'][1]
            dict_milestone['Name'].append(milestone_name)
            dict_milestone['Date'].append(milestone_date)
        milestone_df = pd.DataFrame(dict_milestone)
        milestone_df.to_json(f'{self.proj_dir}\\utils\\memory.json')
        messagebox.showinfo('Info', f'Memory updated!')
